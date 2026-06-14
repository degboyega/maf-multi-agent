"""Generates a branded PowerPoint health report from agent output HTML."""

import io
import logging
import re
from datetime import datetime
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# Brand colours from Eagle.pptx
WRAGBY_RED = RGBColor(0xC0, 0x00, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x40, 0x40, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY_BG = RGBColor(0xF2, 0xF2, 0xF2)
STATUS_COLORS = {
    "critical": RGBColor(0xC0, 0x00, 0x00),
    "advisory": RGBColor(0xED, 0x7D, 0x31),
    "normal": RGBColor(0x70, 0xAD, 0x47),
}
STATUS_LABELS = {
    "critical": "● CRITICAL",
    "advisory": "● ADVISORY",
    "normal": "● NORMAL",
}

TITLE_FONT = "Gotham"
BODY_FONT = "Calibri"
TAGLINE = "Work Smart, Achieve More"

TEMPLATE_PATH = Path(__file__).parent.parent.parent / "assets" / "report_template.pptx"

# Slide dimensions (13.33" x 7.50")
W = 13.33
H = 7.50


# ── helpers ──────────────────────────────────────────────────────────────────

def _remove_all_slides(prs: Presentation) -> None:
    slide_id_list = prs.part._element.find(
        ".//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst"
    )
    for sld_id in list(slide_id_list):
        rId = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        slide_id_list.remove(sld_id)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def _add_slide(prs: Presentation, layout_name: str = "Custom Layout"):
    layout = next(
        (l for l in prs.slide_layouts if l.name == layout_name),
        prs.slide_layouts[2],
    )
    return prs.slides.add_slide(layout)


def _title_box(slide, text: str):
    tb = slide.shapes.add_textbox(
        Inches(0.51), Inches(0.53), Inches(9.52), Inches(0.72)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = TITLE_FONT
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WRAGBY_RED
    return tb


def _body_box(slide, lines: list[str], top: float = 1.41, height: float = 4.8):
    tb = slide.shapes.add_textbox(
        Inches(0.65), Inches(top), Inches(12.2), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if line == "":
            continue
        # Bold lines that start with ★ or are section headers
        is_header = line.startswith("★") or (line.endswith(":") and len(line) < 40)
        run = p.add_run()
        run.text = line
        run.font.name = BODY_FONT
        run.font.size = Pt(18 if is_header else 16)
        run.font.bold = is_header
        run.font.color.rgb = WRAGBY_RED if is_header else BLACK
    return tb


def _status_badge(slide, status: str, left: float, top: float):
    color = STATUS_COLORS.get(status.lower(), GREY)
    label = STATUS_LABELS.get(status.lower(), f"● {status.upper()}")
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(3.5), Inches(0.45))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.name = BODY_FONT
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = color
    return tb


def _tagline(slide):
    tb = slide.shapes.add_textbox(
        Inches(W - 2.0), Inches(H - 0.38), Inches(1.9), Inches(0.34)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = TAGLINE
    run.font.name = BODY_FONT
    run.font.size = Pt(9)
    run.font.color.rgb = GREY
    return tb


def _divider(slide, top: float = 1.25):
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.LINE
        Inches(0.51), Inches(top), Inches(12.3), Inches(0),
    )
    line.line.color.rgb = WRAGBY_RED
    line.line.width = Pt(1.5)
    return line


# ── HTML parser ───────────────────────────────────────────────────────────────

def _strip_html(html: str) -> str:
    text = re.sub(r"<br\s*/?>", "\n", html, flags=re.I)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"&amp;", "&", text)
    text = re.sub(r"&lt;", "<", text)
    text = re.sub(r"&gt;", ">", text)
    text = re.sub(r" {2,}", " ", text)
    return text.strip()


def _detect_status(text: str) -> str:
    t = text.lower()
    if "critical" in t:
        return "critical"
    if "advisory" in t or "warning" in t:
        return "advisory"
    if "normal" in t:
        return "normal"
    return "unknown"


def _extract_asset_blocks(plain_text: str) -> list[dict]:
    """Extract per-asset data from the plain-text report."""
    assets = []
    # Split on asset ID headings like "COMP-001" or "## COMP-001"
    pattern = re.compile(r"(COMP[-\s]?\d{3})", re.I)
    parts = pattern.split(plain_text)
    i = 1
    while i < len(parts) - 1:
        asset_id = parts[i].strip().upper().replace(" ", "-")
        block = parts[i + 1] if i + 1 < len(parts) else ""
        i += 2

        # Deduplicate — skip if we already extracted this asset
        if any(a["id"] == asset_id for a in assets):
            continue

        status = _detect_status(block[:200])
        location = ""
        loc_match = re.search(r"(Warri|Eket|Port Harcourt|Lagos|Bonny|Brass|Niger Delta|Akwa Ibom|Delta State|[A-Z][a-z]+,\s*[A-Z][a-z]+)", block)
        if loc_match:
            location = loc_match.group(1)

        # Extract metric lines (lines with numbers + units)
        metrics = []
        for line in block.split("\n"):
            line = line.strip()
            if re.search(r"\d+[°%\s]*(F|psi|rpm|in/s|bar|°C|kPa|kW)", line):
                cleaned = re.sub(r"[|*_]{1,3}", "", line).strip()
                if cleaned and len(cleaned) < 120:
                    metrics.append(cleaned)

        # Extract root cause
        root_cause = ""
        rc_match = re.search(
            r"(root\s*cause|hypothesis|likely\s*cause)[:\s]+([^\n.]{10,200})",
            block, re.I
        )
        if rc_match:
            root_cause = rc_match.group(2).strip()

        # Extract recommendation
        recommendation = ""
        rec_match = re.search(
            r"(recommend(?:ed)?\s*action|action\s*item|next\s*step)[:\s]+([^\n.]{10,250})",
            block, re.I
        )
        if rec_match:
            recommendation = rec_match.group(2).strip()

        assets.append({
            "id": asset_id,
            "location": location,
            "status": status,
            "metrics": metrics[:6],
            "root_cause": root_cause,
            "recommendation": recommendation,
            "raw": block[:800],
        })

    return assets


def _extract_recommendations(plain_text: str) -> list[str]:
    recs = []
    in_section = False
    for line in plain_text.split("\n"):
        stripped = line.strip()
        if re.search(r"recommendation|action\s*item|next\s*step|priority", stripped, re.I):
            in_section = True
            continue
        if in_section:
            if stripped and not stripped.startswith("#"):
                cleaned = re.sub(r"^[-•*\d.]+\s*", "", stripped)
                if len(cleaned) > 10:
                    recs.append(cleaned[:150])
            if len(recs) >= 6:
                break
    return recs


# ── slide builders ────────────────────────────────────────────────────────────

def _build_title_slide(prs: Presentation, report_date: str):
    slide = _add_slide(prs, "Title Slide")
    # Company name
    tb = slide.shapes.add_textbox(Inches(0.65), Inches(1.91), Inches(10.0), Inches(0.77))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "Wragby Business Solutions & Technologies Limited"
    run.font.name = BODY_FONT
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WRAGBY_RED

    # Main title
    tb2 = slide.shapes.add_textbox(Inches(0.55), Inches(3.0), Inches(9.0), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    run2 = tf2.paragraphs[0].add_run()
    run2.text = "Weekly Compressor Health Report"
    run2.font.name = TITLE_FONT
    run2.font.size = Pt(40)
    run2.font.bold = True
    run2.font.color.rgb = BLACK

    p2 = tf2.add_paragraph()
    run3 = p2.add_run()
    run3.text = f"Generated by MAF  ·  {report_date}"
    run3.font.name = BODY_FONT
    run3.font.size = Pt(18)
    run3.font.color.rgb = GREY

    _tagline(slide)


def _build_exec_summary(prs: Presentation, assets: list[dict]):
    slide = _add_slide(prs)
    _title_box(slide, "Executive Summary")
    _divider(slide)

    lines = ["★ Asset Status Overview", ""]
    for a in assets:
        loc = f"  ·  {a['location']}" if a['location'] else ""
        status_label = a['status'].capitalize()
        lines.append(f"{a['id']}{loc}  —  {status_label}")

    lines += ["", "★ Key Findings"]
    for a in assets:
        if a['recommendation']:
            lines.append(f"• {a['id']}: {a['recommendation'][:120]}")
        elif a['status'] in ("advisory", "critical"):
            lines.append(f"• {a['id']}: Requires attention — {a['status'].upper()}")

    _body_box(slide, lines)
    _tagline(slide)


def _build_asset_slide(prs: Presentation, asset: dict):
    slide = _add_slide(prs)
    title = f"{asset['id']}  ·  {asset['location']}" if asset['location'] else asset['id']
    _title_box(slide, title)
    _divider(slide)
    _status_badge(slide, asset['status'], left=10.0, top=0.45)

    lines = []
    if asset['metrics']:
        lines.append("★ Sensor Readings")
        for m in asset['metrics']:
            lines.append(f"• {m}")
        lines.append("")

    if asset['root_cause']:
        lines.append("★ Root Cause Analysis")
        lines.append(asset['root_cause'][:200])
        lines.append("")

    if asset['recommendation']:
        lines.append("★ Recommended Action")
        lines.append(asset['recommendation'][:200])

    if not lines:
        lines = [asset['raw'][:500]]

    _body_box(slide, lines)
    _tagline(slide)


def _build_recommendations_slide(prs: Presentation, assets: list[dict], extra_recs: list[str]):
    slide = _add_slide(prs)
    _title_box(slide, "Recommended Actions")
    _divider(slide)

    lines = []
    priority = 1
    for a in assets:
        if a['recommendation']:
            lines.append(f"Priority {priority}  —  {a['id']}")
            lines.append(f"    {a['recommendation'][:180]}")
            lines.append("")
            priority += 1

    for rec in extra_recs[:3]:
        if not any(rec[:40] in a['recommendation'][:40] for a in assets):
            lines.append(f"Priority {priority}  —  {rec[:180]}")
            lines.append("")
            priority += 1

    if not lines:
        lines = ["No critical actions identified.", "Continue standard monitoring schedule."]

    _body_box(slide, lines)
    _tagline(slide)


# ── public API ────────────────────────────────────────────────────────────────

def build_health_report_pptx(body_html: str, subject: str = "") -> Optional[bytes]:
    """Generate a branded PPTX from the email body HTML. Returns bytes or None on failure."""
    try:
        if not TEMPLATE_PATH.exists():
            logger.warning("⚠️ PPTX template not found at %s — skipping attachment", TEMPLATE_PATH)
            return None

        plain = _strip_html(body_html)
        assets = _extract_asset_blocks(plain)
        extra_recs = _extract_recommendations(plain)
        report_date = datetime.now().strftime("%d %B %Y")

        prs = Presentation(str(TEMPLATE_PATH))
        _remove_all_slides(prs)

        _build_title_slide(prs, report_date)

        if assets:
            _build_exec_summary(prs, assets)
            for asset in assets:
                _build_asset_slide(prs, asset)

        _build_recommendations_slide(prs, assets, extra_recs)

        buf = io.BytesIO()
        prs.save(buf)
        result = buf.getvalue()
        logger.info("📊 PPTX built: %d slides, %.0f KB", len(prs.slides), len(result) / 1024)
        return result

    except Exception as exc:
        logger.error("❌ PPTX build failed: %s", exc, exc_info=True)
        return None
